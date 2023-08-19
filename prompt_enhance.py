import os
import csv
import re
import fitz
import collections
import collections.abc
from pptx import Presentation

def collect_user_details():
    name = input("Please enter your name: ")
    email = input("Please enter your email: ")
    age = int(input("Please enter your age: "))
    location = input("Please enter your location: ")
    learning_pace = input("How would you describe your current proficiency (beginner/intermediate/advanced)? ")
    self_assessment = int(input("On a scale of 1 to 10, how well do you think you know the topic? "))
    return name, email, age, location, learning_pace, self_assessment

def save_to_csv(user_details, csv_file):
    with open(csv_file, "a", newline="") as csvfile:
        csv_writer = csv.writer(csvfile)
        csv_writer.writerow(user_details)

def main():
    custom_path = input("Enter the custom path for the CSV file: ")
    csv_file = os.path.join(custom_path, "user_details.csv")
    name, email, age, location, learning_pace, self_assessment = collect_user_details()
    user_details = [name, email, age, location, learning_pace, self_assessment]
    save_to_csv(user_details, csv_file)
    print("User details collected and stored in", csv_file)
    save_to_csv(user_details, csv_file)
    print("User details updated and stored in", csv_file)
    return learning_pace

def pdf_to_text():
    pdf_path = input("Enter path of PDF: ")
    start_page = int(input("Enter start page: ")) - 1
    end_page = int(input("Enter end page: ")) - 1
    doc = fitz.open(pdf_path)
    if start_page < 0 or end_page >= doc.page_count or start_page > end_page:
        print(f"Invalid page range. Total pages: {doc.page_count}")
        doc.close()
        return
    for page_num in range(start_page, end_page + 1):
        page = doc.load_page(page_num)
        text = page.get_text("text")
    doc.close()
    return text

def quiz(prompt):
    pattern = re.compile(r'Answer:\s*([A-D])\.\s')
    matches = pattern.findall(prompt)
    questions = re.split(r'Answer:\s*[A-D]\.\s', prompt)
    questions = [q.strip() for q in questions if q.strip()]
    for i, question in enumerate(questions):
        print(question)
        if i < len(matches):
            correct_answer = matches[i]
            user_answer = input("Enter your answer: ").upper()
            if user_answer == correct_answer:
                print("Correct!\n")
            else:
                print(f"Incorrect! The correct answer is {correct_answer}\n")

def ppt_to_text():
    ppt_path = input("Enter path of PowerPoint file: ")
    presentation = Presentation(ppt_path)
    text_content = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    text=" ".join(text_content)
    return text
